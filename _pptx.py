import pptx
from pptx import Presentation
import os
from random import seed
from random import randint


class ExtractPPTX:
    def __init__(self, file_name):
        self.pptx = Presentation(file_name)
        self.raw_text = []

    def function(self):

        slide_counter = 1

        current = 1

        for slide in self.pptx.slides:

            slide_number = str(slide_counter)

            slide_counter += 1

            text_run = []

            ran_int = randint(0, 1000)

            for shape in slide.shapes:

                if hasattr(shape, "image"):
                    path = "images"
                    print(str(shape.image.sha1))
                    file_name = str(shape.image.sha1)

                    # check if path to images exist, make it if not
                    if not os.path.exists(path):
                        os.makedirs(path)

                    file_path = (
                        "images/" + slide_number + "-" + file_name + "image" + ".png"
                    )

                    with open(file_path, "wb+") as fo:
                        fo.write(shape.image.blob)
                        fo.close()

                    text_run.append("[img" + slide_number + "]")

                if not shape.has_text_frame:
                    continue

                check = any("Slide " + slide_number in s for s in text_run)

                current = slide_number

                if check:
                    text_run.append("\n\n")
                if not check:
                    text_run.append("Slide " + slide_number + "\n")

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_run.append(run.text + " \n")

            text_run.append("\n------------------------------\n")
            extracted_text = "".join(text_run)
            self.raw_text.append(extracted_text)

        print("Finished Converting.....")
        return self.raw_text

    def get_slide_info(self, file_name):
        self.pptx = Presentation(file_name)
        slide = []
        for s in self.pptx.slides:
            slide.append(self.pptx.slides.index(s) + 1)
        return slide
